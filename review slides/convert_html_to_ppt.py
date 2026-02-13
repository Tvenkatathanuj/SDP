"""
Convert HTML slides to PowerPoint by capturing them as high-quality images
This preserves the exact visual appearance of the HTML slides
"""
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os
import subprocess
import sys

def install_playwright_browsers():
    """Install playwright browsers if not already installed"""
    print("Installing Playwright browsers (one-time setup)...")
    try:
        result = subprocess.run(
            [sys.executable, "-m", "playwright", "install", "chromium"],
            capture_output=True,
            text=True,
            timeout=120
        )
        if result.returncode == 0:
            print("✓ Playwright browsers installed successfully")
        else:
            print(f"Warning: {result.stderr}")
    except Exception as e:
        print(f"Note: {e}")

def capture_html_to_image(html_file, output_image, width=1280, height=720):
    """Capture HTML file as PNG image using Playwright"""
    from playwright.sync_api import sync_playwright
    
    abs_path = os.path.abspath(html_file)
    file_url = f"file:///{abs_path.replace(os.sep, '/')}"
    
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={'width': width, 'height': height})
        page.goto(file_url)
        page.wait_for_load_state('networkidle')
        page.screenshot(path=output_image, full_page=False)
        browser.close()

def create_presentation_from_images(image_files, output_pptx):
    """Create PowerPoint presentation from image files"""
    prs = Presentation()
    
    # Set slide size to 16:9 (standard presentation size)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 ratio
    
    for idx, img_file in enumerate(image_files, 1):
        print(f"Adding slide {idx}/{len(image_files)}: {os.path.basename(img_file)}")
        
        # Add blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add image to fill the entire slide
        slide.shapes.add_picture(
            img_file,
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
    
    prs.save(output_pptx)
    print(f"\n✓ Presentation saved as: {output_pptx}")

def main():
    print("="*60)
    print("HTML to PowerPoint Converter")
    print("High-Quality Image-Based Conversion")
    print("="*60)
    print()
    
    # Install playwright browsers (one-time)
    install_playwright_browsers()
    print()
    
    # Get all HTML files
    html_files = [f"s{i}.html" for i in range(1, 17)]
    html_files = [f for f in html_files if os.path.exists(f)]
    
    print(f"Found {len(html_files)} HTML files\n")
    
    # Create temp directory for images
    img_dir = "temp_slide_images"
    if not os.path.exists(img_dir):
        os.makedirs(img_dir)
    
    # Capture each HTML file as image
    image_files = []
    for idx, html_file in enumerate(html_files, 1):
        print(f"[{idx}/{len(html_files)}] Capturing {html_file}...")
        img_file = os.path.join(img_dir, f"slide_{idx:02d}.png")
        
        try:
            capture_html_to_image(html_file, img_file)
            image_files.append(img_file)
            print(f"    ✓ Saved to {img_file}")
        except Exception as e:
            print(f"    ✗ Error: {e}")
    
    print(f"\n{'='*60}")
    print(f"Captured {len(image_files)} slides successfully")
    print(f"{'='*60}\n")
    
    # Create PowerPoint
    if image_files:
        output_pptx = "Parkinsons_PD_Detection_Presentation.pptx"
        print("Creating PowerPoint presentation...")
        create_presentation_from_images(image_files, output_pptx)
        
        print(f"\n{'='*60}")
        print(f"SUCCESS!")
        print(f"{'='*60}")
        print(f"Your presentation is ready: {output_pptx}")
        print(f"Total slides: {len(image_files)}")
        print(f"\nNote: Temp images saved in '{img_dir}/' folder")
        print("You can delete this folder after verifying the presentation.")
    else:
        print("No slides were captured. Please check for errors above.")

if __name__ == "__main__":
    main()
