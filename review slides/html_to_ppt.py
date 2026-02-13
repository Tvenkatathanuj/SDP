"""
Convert HTML slides to PowerPoint presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from bs4 import BeautifulSoup
import os
import re

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def extract_text_content(element):
    """Extract text content from HTML element, preserving structure"""
    if element is None:
        return ""
    return element.get_text(strip=True, separator=' ')

def parse_html_table(table_elem):
    """Parse HTML table and extract data"""
    headers = []
    rows = []
    
    # Extract headers
    thead = table_elem.find('thead')
    if thead:
        header_cells = thead.find_all('th')
        for th in header_cells:
            headers.append(extract_text_content(th))
    
    # Extract rows
    tbody = table_elem.find('tbody')
    if tbody:
        tr_elements = tbody.find_all('tr')
        for tr in tr_elements:
            row_data = []
            cells = tr.find_all('td')
            for td in cells:
                row_data.append(extract_text_content(td))
            if row_data:
                rows.append(row_data)
    
    return headers, rows

def parse_html_slide(html_file):
    """Parse HTML file and extract content"""
    with open(html_file, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')
    
    # Extract title
    title = soup.find('title')
    title_text = title.get_text() if title else ""
    
    # Extract main title/header
    header_text = ""
    header = soup.find(class_=['header-section', 'header-bar'])
    if header:
        title_elem = header.find(class_=['title-text', 'main-title'])
        if title_elem:
            header_text = extract_text_content(title_elem)
    
    # If no header found, try to find main-title in content
    if not header_text:
        main_title = soup.find(class_='main-title')
        if main_title:
            header_text = extract_text_content(main_title)
    
    # Check for tables
    table = soup.find('table')
    table_data = None
    if table:
        headers, rows = parse_html_table(table)
        table_data = {'headers': headers, 'rows': rows}
    
    # Extract bullet points
    bullets = []
    bullet_items = soup.find_all(class_='bullet-item')
    for item in bullet_items:
        bullet_text_elem = item.find(class_='bullet-text')
        if bullet_text_elem:
            text = extract_text_content(bullet_text_elem)
            bullets.append(text)
    
    # If no bullets, try to extract other content
    if not bullets:
        # Look for section-box items
        section_boxes = soup.find_all(class_='section-box')
        for box in section_boxes:
            title_elem = box.find(class_='box-title')
            desc_elem = box.find(class_='box-description')
            if title_elem:
                text = extract_text_content(title_elem)
                if desc_elem:
                    text += ": " + extract_text_content(desc_elem)
                bullets.append(text)
        
        # Look for list items
        if not bullets:
            list_items = soup.find_all('li')
            for li in list_items:
                text = extract_text_content(li)
                if text:
                    bullets.append(text)
        
        # Look for info-card items (for architecture diagrams)
        if not bullets:
            info_cards = soup.find_all(class_='info-card')
            for card in info_cards:
                card_title = card.find(class_='card-title')
                card_content = card.find(class_='card-content')
                if card_title:
                    text = extract_text_content(card_title)
                    if card_content:
                        text += ": " + extract_text_content(card_content)
                    bullets.append(text)
        
        # Look for module-box items (for architecture flow)
        if not bullets:
            module_boxes = soup.find_all(class_='module-box')
            for box in module_boxes:
                module_title = box.find(class_='module-title')
                module_desc = box.find(class_='module-desc')
                if module_title:
                    text = extract_text_content(module_title)
                    if module_desc:
                        text += ": " + extract_text_content(module_desc)
                    bullets.append(text)
    
    # Extract subtitle or additional info
    subtitle = ""
    subtitle_elem = soup.find(class_='subtitle')
    if subtitle_elem:
        subtitle = extract_text_content(subtitle_elem)
    
    # For cover slide, extract team info
    team_info = []
    team_section = soup.find(class_='team-section')
    if team_section:
        members = team_section.find_all(class_='team-member')
        for member in members:
            team_info.append(extract_text_content(member))
    
    guide_info = ""
    guide_section = soup.find(class_='guide-section')
    if guide_section:
        guide_name = guide_section.find(class_='guide-name')
        if guide_name:
            guide_info = "Project Guide: " + extract_text_content(guide_name)
    
    return {
        'title': title_text,
        'header': header_text,
        'subtitle': subtitle,
        'bullets': bullets,
        'team_info': team_info,
        'guide_info': guide_info,
        'table_data': table_data
    }

def add_title_slide(prs, data):
    """Add title slide to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add blue header bar at top
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(30, 58, 138)  # #1e3a8a
    header_shape.line.color.rgb = RGBColor(30, 58, 138)
    
    # Add main title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = data['header']
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)
    
    # Add subtitle
    if data['subtitle']:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.8),
            Inches(8), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = data['subtitle']
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(59, 130, 246)
    
    # Add team info
    if data['team_info']:
        y_pos = 4.5
        team_box = slide.shapes.add_textbox(
            Inches(1), Inches(y_pos),
            Inches(4), Inches(1.5)
        )
        team_frame = team_box.text_frame
        p = team_frame.paragraphs[0]
        p.text = "Presented By"
        p.font.size = Pt(16)
        p.font.bold = True
        for member in data['team_info']:
            p = team_frame.add_paragraph()
            p.text = member
            p.font.size = Pt(12)
    
    # Add guide info
    if data['guide_info']:
        guide_box = slide.shapes.add_textbox(
            Inches(5.5), Inches(4.5),
            Inches(4), Inches(1)
        )
        guide_frame = guide_box.text_frame
        p = guide_frame.paragraphs[0]
        p.text = data['guide_info']
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(14)
        p.font.bold = True
    
    return slide

def add_table_to_slide(slide, table_data, left, top, width, height):
    """Add table to slide"""
    if not table_data or not table_data['headers']:
        return
    
    headers = table_data['headers']
    rows = table_data['rows']
    
    # Create table shape
    rows_count = len(rows) + 1  # +1 for header
    cols_count = len(headers)
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
    table = table_shape.table
    
    # Set column widths evenly
    col_width = width / cols_count
    for i in range(cols_count):
        table.columns[i].width = int(col_width)
    
    # Add headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 58, 138)
        
        # Format header text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Add data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            if col_idx < cols_count:  # Ensure we don't exceed column count
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = cell_data
                
                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(248, 250, 252)
                
                # Format cell text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.alignment = PP_ALIGN.CENTER if col_idx in [0, len(row_data)-1] else PP_ALIGN.LEFT
    
    return table_shape

def add_content_slide(prs, data):
    """Add content slide with bullets or table to presentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add blue header bar
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = RGBColor(30, 58, 138)
    header_shape.line.color.rgb = RGBColor(30, 58, 138)
    
    # Add title text in header
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = data['header'] if data['header'] else data['title']
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Check if slide has table data
    if data.get('table_data'):
        # Add table
        add_table_to_slide(
            slide, 
            data['table_data'],
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5)
        )
    else:
        # Add content area with bullets
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        # Add bullets
        for i, bullet in enumerate(data['bullets']):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = Pt(14)
            p.space_before = Pt(4)
            p.space_after = Pt(4)
    
    return slide

def create_presentation(html_files, output_file):
    """Create PowerPoint presentation from HTML files"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for html_file in html_files:
        print(f"Processing {html_file}...")
        data = parse_html_slide(html_file)
        
        # First slide is title slide
        if 's1.html' in html_file:
            add_title_slide(prs, data)
        else:
            add_content_slide(prs, data)
    
    prs.save(output_file)
    print(f"\nPresentation saved as: {output_file}")

if __name__ == "__main__":
    # Get all HTML files
    html_files = [f"s{i}.html" for i in range(1, 17)]
    
    # Filter only existing files
    html_files = [f for f in html_files if os.path.exists(f)]
    
    print(f"Found {len(html_files)} HTML files")
    
    # Create presentation
    output_file = "PD_Detection_Slides.pptx"
    create_presentation(html_files, output_file)
